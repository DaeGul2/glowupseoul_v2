"""
Generate docs/admin_guidebook.pptx — friendly, non-developer-facing guide.
Run:   python scripts/build_admin_guidebook.py

Layout rules (do NOT break):
  - 슬라이드 영역 (13.333 x 7.5 inches) 안에 모든 요소 들어가야 함
  - 텍스트박스는 word_wrap + auto_size 로 늘어남
  - 한 줄에 안 들어갈 텍스트는 height 충분히 큰 박스 + spacing 보수적
  - 한국어 본문 12pt, sub 10pt, headline 36pt 기준
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

# ---------------------------------------------------------------
# Brand colors
# ---------------------------------------------------------------
BG          = RGBColor(0xFA, 0xFA, 0xF7)
BG_SOFT     = RGBColor(0xF2, 0xEF, 0xE8)
TEXT        = RGBColor(0x2A, 0x2A, 0x2A)
TEXT_SOFT   = RGBColor(0x5A, 0x5A, 0x5A)
TEXT_MUTED  = RGBColor(0x8A, 0x8A, 0x8A)
GOLD        = RGBColor(0xC9, 0xA0, 0x63)
GOLD_DARK   = RGBColor(0x9A, 0x77, 0x44)
GOLD_LIGHT  = RGBColor(0xE8, 0xD4, 0xB2)
WARN        = RGBColor(0xC4, 0x55, 0x4D)
WARN_BG     = RGBColor(0xFA, 0xEC, 0xEA)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LINE        = RGBColor(0xE1, 0xDD, 0xD3)

KO_FONT = "맑은 고딕"
EN_FONT = "Calibri"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ---------------------------------------------------------------
# Drawing helpers
# ---------------------------------------------------------------
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, text, *, size=12, bold=False, italic=False,
             color=TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=KO_FONT, line_spacing=1.3, auto_grow=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    if auto_grow:
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
    return tb

def add_bg(slide):
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, BG)

def add_left_marker(slide):
    add_rect(slide, 0, 0, Inches(0.08), prs.slide_height, GOLD)

def add_page_no(slide, n, total):
    add_text(slide, Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3),
             f"{n} / {total}", size=9, color=TEXT_MUTED,
             align=PP_ALIGN.RIGHT, font=EN_FONT, auto_grow=False)

def add_brand_mark(slide):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(7), Inches(0.3),
             "✦  Glow Up Seoul · 운영 가이드북",
             size=9, color=TEXT_MUTED, auto_grow=False)

def add_eyebrow(slide, text, color=GOLD):
    add_text(slide, Inches(0.7), Inches(0.6), Inches(8), Inches(0.35),
             text, size=10, color=color, bold=True, auto_grow=False)

def add_meta_right(slide, text):
    add_text(slide, Inches(8.5), Inches(0.6), Inches(4.2), Inches(0.35),
             text, size=10, color=TEXT_MUTED, align=PP_ALIGN.RIGHT, auto_grow=False)

def add_h1(slide, text, color=TEXT):
    add_text(slide, Inches(0.7), Inches(1.0), Inches(12), Inches(1.0),
             text, size=36, bold=True, color=color, line_spacing=1.1, auto_grow=False)

def add_subtitle(slide, text, y=Inches(2.0)):
    add_text(slide, Inches(0.7), y, Inches(12), Inches(0.6),
             text, size=15, italic=True, color=TEXT_SOFT,
             line_spacing=1.4, auto_grow=False)

def add_callout(slide, x, y, w, title, body, *, color=WARN, bg=WARN_BG, icon="⚠"):
    """Self-sizing callout — height is fixed at 1.0 inch (room for 2-line body)."""
    h = Inches(1.0)
    add_rect(slide, x, y, w, h, bg)
    add_rect(slide, x, y, Inches(0.06), h, color)
    add_text(slide, x + Inches(0.3), y + Inches(0.15), Inches(0.5), Inches(0.4),
             icon, size=18, color=color, bold=True, auto_grow=False)
    add_text(slide, x + Inches(0.85), y + Inches(0.15), w - Inches(1.1), Inches(0.35),
             title, size=12, bold=True, color=TEXT, auto_grow=False)
    add_text(slide, x + Inches(0.85), y + Inches(0.5), w - Inches(1.1), Inches(0.45),
             body, size=10.5, color=TEXT_SOFT, line_spacing=1.45, auto_grow=False)

# Slide registry
SLIDES = []
def reg(fn):
    SLIDES.append(fn); return fn

# ===============================================================
# 01 COVER
# ===============================================================
@reg
def s01(n, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_rect(s, 0, Inches(6.4), prs.slide_width, Inches(0.08), GOLD)
    add_text(s, Inches(1.0), Inches(0.9), Inches(4), Inches(0.4),
             "✦  GLOW UP SEOUL", size=11, color=GOLD, bold=True, auto_grow=False)
    add_text(s, Inches(1.0), Inches(2.0), Inches(11), Inches(1.5),
             "운영 가이드북", size=72, bold=True, color=TEXT,
             line_spacing=1.0, auto_grow=False)
    add_text(s, Inches(1.0), Inches(3.5), Inches(11), Inches(0.7),
             "Admin 화면 사용법 · 비개발자용",
             size=24, italic=True, color=TEXT_SOFT, auto_grow=False)
    add_text(s, Inches(1.0), Inches(5.3), Inches(11), Inches(0.4),
             "처음 한 번 정독한 뒤,", size=13, color=TEXT_SOFT, auto_grow=False)
    add_text(s, Inches(1.0), Inches(5.7), Inches(11), Inches(0.4),
             "필요할 때마다 해당 챕터만 다시 보시면 됩니다.",
             size=13, color=TEXT_SOFT, auto_grow=False)
    add_text(s, Inches(1.0), Inches(6.8), Inches(6), Inches(0.3),
             "v1 · 2026-05-14", size=11, color=TEXT_MUTED, auto_grow=False)
    add_text(s, Inches(7.0), Inches(6.8), Inches(5.3), Inches(0.3),
             "문의 · glowupinseoul@gmail.com",
             size=11, color=TEXT_MUTED, align=PP_ALIGN.RIGHT, auto_grow=False)

# ===============================================================
# 02 이 가이드는 누구를 위한 것?
# ===============================================================
@reg
def s02(n, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s); add_left_marker(s); add_brand_mark(s); add_page_no(s, n, total)
    add_eyebrow(s, "◇  시작하기 전에")
    add_h1(s, "이 가이드는 누구를 위한 것인가요?")
    add_subtitle(s, "글로업 서울의 운영을 도와주시는 분 — 즉 당신을 위한 것입니다.",
                 y=Inches(2.2))

    add_text(s, Inches(0.7), Inches(3.1), Inches(12), Inches(0.4),
             "이 가이드는 다음을 가정합니다", size=14, color=TEXT, bold=True,
             auto_grow=False)

    items = [
        "프로그래밍을 모르셔도 됩니다. 코드 한 줄도 안 나옵니다.",
        "관리자 페이지(Admin)에 들어가는 방법은 알고 있다고 봅니다.",
        "병원·시술·환자 후기 등을 직접 입력하실 분입니다.",
        "막혔을 때 누구에게 물어볼지는 알고 계십니다 (마지막 슬라이드 참고).",
    ]
    cy = Inches(3.7)
    for txt in items:
        add_text(s, Inches(0.85), cy, Inches(0.3), Inches(0.4),
                 "◇", size=14, color=GOLD, bold=True, auto_grow=False)
        add_text(s, Inches(1.2), cy, Inches(11.2), Inches(0.5),
                 txt, size=13, color=TEXT, auto_grow=False)
        cy += Inches(0.55)

    add_callout(s, Inches(0.7), Inches(6.1), Inches(11.9),
                "딱 한 가지 약속",
                "이 가이드의 순서대로 따라가시면 사이트가 정상 작동합니다. 순서를 건너뛰면 사이트가 텅 비거나 잘못된 정보가 보일 수 있어요.",
                color=GOLD, bg=GOLD_LIGHT, icon="✦")

# ===============================================================
# 03 Admin 화면 구조
# ===============================================================
@reg
def s03(n, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s); add_left_marker(s); add_brand_mark(s); add_page_no(s, n, total)
    add_eyebrow(s, "◇  먼저 익숙해지기")
    add_h1(s, "Admin 화면의 구조")
    add_subtitle(s, "왼쪽 메뉴 9개. 각 메뉴는 한 가지 데이터를 다룹니다.",
                 y=Inches(2.1))

    # Left: sidebar mockup
    sb_x, sb_y, sb_w, sb_h = Inches(0.7), Inches(3.0), Inches(3.6), Inches(4.0)
    add_rect(s, sb_x, sb_y, sb_w, sb_h, WHITE, line=LINE)
    add_text(s, sb_x + Inches(0.25), sb_y + Inches(0.15), sb_w - Inches(0.5), Inches(0.3),
             "사이드바", size=11, bold=True, color=GOLD, auto_grow=False)
    menu = [
        ("📊", "대시보드"),
        ("📥", "파트너 신청서"),
        ("🏥", "병원"),
        ("✦", "시술"),
        ("◇", "카테고리"),
        ("✿", "고민"),
        ("👤", "의사"),
        ("📸", "전후사진"),
        ("📰", "실시간 피드"),
        ("🔗", "고민↔시술 매트릭스"),
    ]
    cy = sb_y + Inches(0.55)
    for ic, label in menu:
        add_text(s, sb_x + Inches(0.3), cy, sb_w - Inches(0.6), Inches(0.32),
                 f"{ic}  {label}", size=11, color=TEXT, auto_grow=False)
        cy += Inches(0.32)

    # Right: explanation cards
    rx = Inches(4.7); ry = Inches(3.0); rw = Inches(8.0)
    add_text(s, rx, ry, rw, Inches(0.4),
             "메뉴별 사용 빈도", size=15, bold=True, color=TEXT, auto_grow=False)

    cards = [
        ("자주 가는 곳 (운영의 70%)",
         "병원 · 시술 · 고민↔시술 매트릭스"),
        ("처음 채우고 잘 안 만지는 곳",
         "카테고리 · 의사 · 전후사진"),
        ("매일 보는 곳",
         "대시보드 (전체 카운트) · 파트너 신청서 (새 신청)"),
    ]
    cy = ry + Inches(0.55)
    for title, body in cards:
        add_rect(s, rx, cy, rw, Inches(1.05), WHITE, line=LINE)
        add_rect(s, rx, cy, Inches(0.06), Inches(1.05), GOLD)
        add_text(s, rx + Inches(0.25), cy + Inches(0.15), rw - Inches(0.4), Inches(0.35),
                 title, size=13, bold=True, color=TEXT, auto_grow=False)
        add_text(s, rx + Inches(0.25), cy + Inches(0.55), rw - Inches(0.4), Inches(0.4),
                 body, size=11, color=TEXT_SOFT, auto_grow=False)
        cy += Inches(1.2)

# ===============================================================
# 04 큰 그림 — 입력 순서
# ===============================================================
@reg
def s04(n, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s); add_left_marker(s); add_brand_mark(s); add_page_no(s, n, total)
    add_eyebrow(s, "◇  큰 그림")
    add_h1(s, "어떤 순서로 채워야 하나요?")
    add_subtitle(s, "데이터는 서로 연결돼 있어서 순서가 중요합니다.",
                 y=Inches(2.1))

    steps = [
        ("①", "카테고리 사진",   "(자동 완료) 사진만 추가"),
        ("②", "고민",          "고객 검색 키워드 (Lifting · Wrinkles 등)"),
        ("③", "시술",          "우리가 다루는 시술 카탈로그"),
        ("④", "매트릭스",       "고민 ↔ 시술 연결 (⭐ 가장 중요)"),
        ("⑤", "병원",          "위 모든 것을 사용해서 등록"),
        ("⑥", "실시간 피드",    "메인 화면 후기 부트스트랩"),
    ]
    # Two columns × 3 rows
    col_w = Inches(5.9); col_gap = Inches(0.3)
    row_h = Inches(1.05); row_gap = Inches(0.2)
    grid_x = Inches(0.7); grid_y = Inches(2.85)
    for i, (num, title, body) in enumerate(steps):
        col = i % 2; row = i // 2
        x = grid_x + (col_w + col_gap) * col
        y = grid_y + (row_h + row_gap) * row
        add_rect(s, x, y, col_w, row_h, WHITE, line=LINE)
        add_rect(s, x, y, Inches(0.08), row_h, GOLD)
        add_text(s, x + Inches(0.25), y + Inches(0.18), Inches(0.7), Inches(0.55),
                 num, size=22, bold=True, color=GOLD, font=EN_FONT, auto_grow=False)
        add_text(s, x + Inches(0.95), y + Inches(0.18), col_w - Inches(1.1), Inches(0.4),
                 title, size=14, bold=True, color=TEXT, auto_grow=False)
        add_text(s, x + Inches(0.95), y + Inches(0.58), col_w - Inches(1.1), Inches(0.4),
                 body, size=11, color=TEXT_SOFT, auto_grow=False)

    add_callout(s, Inches(0.7), Inches(6.4), Inches(11.9),
                "순서를 어기지 마세요",
                "예: 시술이 없는 상태에서 병원에 시술 가격표를 추가하려고 하면, 선택할 시술이 없어서 막힙니다.",
                color=WARN, bg=WARN_BG)

# ===============================================================
# 05 함정 TOP 5
# ===============================================================
@reg
def s05(n, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s); add_left_marker(s); add_brand_mark(s); add_page_no(s, n, total)
    add_eyebrow(s, "⚠  미리 알아두기", color=WARN)
    add_h1(s, "가장 자주 빠지는 함정 5가지")
    add_subtitle(s, "본격 시작 전에 이 다섯 가지만 미리 알아두세요.",
                 y=Inches(2.1))

    traps = [
        ("'계약 상태' 가 '활성' 이 아니에요",
         "기본값이 '대기 중'이라 사이트에 안 보여요. 저장 후 꼭 '활성'으로."),
        ("URL 식별자 한 번 정하면 절대 바꾸지 마세요",
         "외부 링크와 사진 폴더가 이 이름으로 묶여 있어 바꾸면 다 깨집니다."),
        ("전후사진 — 환자 서면 동의 없이 절대 X",
         "한국 개인정보보호법(PIPA) 위반. '동의 받음' 체크 꼭 확인."),
        ("사진을 비워두지 마세요",
         "썸네일 비어있으면 사이트 카드가 회색 박스로 표시됩니다."),
        ("입력 순서를 어기지 마세요",
         "다음 슬라이드부터의 순서 그대로 따라가시면 절대 안 막힙니다."),
    ]
    cy = Inches(2.85)
    row_h = Inches(0.82)
    for i, (t, b) in enumerate(traps, 1):
        add_rect(s, Inches(0.7), cy, Inches(11.9), row_h, WHITE, line=LINE)
        # number badge
        add_rect(s, Inches(0.9), cy + Inches(0.18), Inches(0.45), Inches(0.45), WARN)
        add_text(s, Inches(0.9), cy + Inches(0.22), Inches(0.45), Inches(0.4),
                 str(i), size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font=EN_FONT, auto_grow=False)
        # title
        add_text(s, Inches(1.55), cy + Inches(0.12), Inches(11.0), Inches(0.35),
                 t, size=13, bold=True, color=TEXT, auto_grow=False)
        # body
        add_text(s, Inches(1.55), cy + Inches(0.45), Inches(11.0), Inches(0.32),
                 b, size=10.5, color=TEXT_SOFT, auto_grow=False)
        cy += row_h + Inches(0.04)

# ===============================================================
# Step helper
# ===============================================================
def step_slide(n, total, *, step_label, step_title, est_time,
               where, click_path, fields, why, sample=None, warn=None):
    s = prs.slides.add_slide(BLANK)
    add_bg(s); add_left_marker(s); add_brand_mark(s); add_page_no(s, n, total)
    add_eyebrow(s, step_label)
    add_meta_right(s, f"예상 시간 · {est_time}")
    add_h1(s, step_title)

    # Where + Click strip
    add_rect(s, Inches(0.7), Inches(2.05), Inches(11.9), Inches(0.75), BG_SOFT)
    add_text(s, Inches(0.95), Inches(2.15), Inches(1.3), Inches(0.3),
             "📍  어디로", size=10, bold=True, color=GOLD, auto_grow=False)
    add_text(s, Inches(2.35), Inches(2.13), Inches(10.1), Inches(0.32),
             where, size=12, color=TEXT, auto_grow=False)
    add_text(s, Inches(0.95), Inches(2.48), Inches(1.3), Inches(0.3),
             "👆  순서", size=10, bold=True, color=GOLD, auto_grow=False)
    add_text(s, Inches(2.35), Inches(2.48), Inches(10.1), Inches(0.32),
             click_path, size=11, color=TEXT_SOFT, italic=True, auto_grow=False)

    # Two columns
    lx, ly = Inches(0.7), Inches(3.05)
    lw = Inches(7.0)
    rx, ry = Inches(7.95), Inches(3.05)
    rw = Inches(4.7)

    # FIELDS column (left)
    add_text(s, lx, ly, lw, Inches(0.35),
             "✏️  채워야 할 것", size=13, bold=True, color=TEXT, auto_grow=False)
    field_gap = Inches(0.55)
    cy = ly + Inches(0.45)
    for label, hint, required in fields:
        # bullet
        add_text(s, lx, cy, Inches(0.3), Inches(0.3),
                 "●" if required else "○",
                 size=13, color=(WARN if required else GOLD), auto_grow=False)
        # label
        add_text(s, lx + Inches(0.3), cy, lw - Inches(0.3), Inches(0.3),
                 label, size=11.5, bold=True, color=TEXT, auto_grow=False)
        # hint
        add_text(s, lx + Inches(0.3), cy + Inches(0.27), lw - Inches(0.3), Inches(0.26),
                 hint, size=9.5, color=TEXT_MUTED, line_spacing=1.3, auto_grow=False)
        cy += field_gap

    # WHY column (right)
    add_text(s, rx, ry, rw, Inches(0.35),
             "💡  왜 채우나요?", size=13, bold=True, color=TEXT, auto_grow=False)
    add_text(s, rx, ry + Inches(0.4), rw, Inches(2.0),
             why, size=11, color=TEXT_SOFT, line_spacing=1.5, auto_grow=False)

    if sample:
        sy = ry + Inches(2.5)
        add_rect(s, rx, sy, rw, Inches(1.05), WHITE, line=LINE)
        add_rect(s, rx, sy, Inches(0.05), Inches(1.05), GOLD)
        add_text(s, rx + Inches(0.2), sy + Inches(0.12), rw - Inches(0.3), Inches(0.3),
                 f"예시 · {sample[0]}", size=11, bold=True, color=TEXT, auto_grow=False)
        add_text(s, rx + Inches(0.2), sy + Inches(0.45), rw - Inches(0.3), Inches(0.55),
                 sample[1], size=10, color=TEXT_SOFT,
                 line_spacing=1.4, auto_grow=False)

    if warn:
        add_callout(s, Inches(0.7), Inches(6.35), Inches(11.9),
                    warn[0], warn[1],
                    color=warn[2] if len(warn) > 2 else WARN,
                    bg=warn[3] if len(warn) > 3 else WARN_BG,
                    icon=warn[4] if len(warn) > 4 else "⚠")
    return s

# ===============================================================
# Usage map helper — "이 사진은 사이트 어디에 나오나"
# ===============================================================
def usage_slide(n, total, *, label, title, summary, slots, pages):
    """
    slots: list of (id, name, size, note)
    pages: list of (page_name, route, items) where items = list of (slot_id, location_desc)
    """
    s = prs.slides.add_slide(BLANK)
    add_bg(s); add_left_marker(s); add_brand_mark(s); add_page_no(s, n, total)
    add_eyebrow(s, label)
    add_h1(s, title)
    add_subtitle(s, summary, y=Inches(2.1))

    # Left column — photo slots
    lx, ly = Inches(0.7), Inches(2.95)
    lw = Inches(5.5)
    add_text(s, lx, ly, lw, Inches(0.35),
             "📷  사진 슬롯 (admin 입력)", size=12, bold=True, color=TEXT, auto_grow=False)
    cy = ly + Inches(0.45)
    for sid, sname, ssize, snote in slots:
        # card
        add_rect(s, lx, cy, lw, Inches(0.8), WHITE, line=LINE)
        add_rect(s, lx, cy, Inches(0.05), Inches(0.8), GOLD)
        # number badge
        add_rect(s, lx + Inches(0.2), cy + Inches(0.18), Inches(0.35), Inches(0.35), GOLD)
        add_text(s, lx + Inches(0.2), cy + Inches(0.2), Inches(0.35), Inches(0.32),
                 sid, size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font=EN_FONT, auto_grow=False)
        # name + size
        add_text(s, lx + Inches(0.7), cy + Inches(0.12), lw - Inches(0.8), Inches(0.3),
                 sname, size=12, bold=True, color=TEXT, auto_grow=False)
        add_text(s, lx + Inches(0.7), cy + Inches(0.38), lw - Inches(0.8), Inches(0.22),
                 ssize, size=9.5, color=GOLD, italic=True, auto_grow=False)
        # hint
        add_text(s, lx + Inches(0.7), cy + Inches(0.58), lw - Inches(0.8), Inches(0.22),
                 snote, size=9.5, color=TEXT_MUTED, auto_grow=False)
        cy += Inches(0.92)

    # Right column — site usage
    rx, ry = Inches(6.6), Inches(2.95)
    rw = Inches(6.0)
    add_text(s, rx, ry, rw, Inches(0.35),
             "🌐  사이트의 어디에 보여요?", size=12, bold=True, color=TEXT, auto_grow=False)
    cy = ry + Inches(0.45)
    for page_name, route, items in pages:
        # page header
        add_rect(s, rx, cy, rw, Inches(0.42), BG_SOFT)
        add_text(s, rx + Inches(0.2), cy + Inches(0.08), rw - Inches(0.4), Inches(0.28),
                 page_name, size=11, bold=True, color=TEXT, auto_grow=False)
        add_text(s, rx + Inches(0.2), cy + Inches(0.27), rw - Inches(0.4), Inches(0.18),
                 route, size=9, color=TEXT_MUTED, italic=True,
                 font=EN_FONT, auto_grow=False)
        cy += Inches(0.5)
        # items
        for sid, desc in items:
            # slot id badge
            add_rect(s, rx + Inches(0.1), cy + Inches(0.04), Inches(0.28), Inches(0.28), GOLD_LIGHT)
            add_text(s, rx + Inches(0.1), cy + Inches(0.06), Inches(0.28), Inches(0.25),
                     sid, size=11, bold=True, color=GOLD_DARK,
                     align=PP_ALIGN.CENTER, font=EN_FONT, auto_grow=False)
            # description
            add_text(s, rx + Inches(0.5), cy + Inches(0.04), rw - Inches(0.6), Inches(0.32),
                     desc, size=10.5, color=TEXT_SOFT, auto_grow=False)
            cy += Inches(0.38)
        cy += Inches(0.1)
    return s

# ===============================================================
# Photo usage slide — annotated screenshot of real site
# ===============================================================
SCREEN_DIR = os.path.normpath(os.path.join(os.path.dirname(__file__), '..', 'docs', 'screenshots', 'annotated'))

def photo_usage_slide(n, total, *, label, title, lede,
                      slots, images=None, fallback_note=None):
    """
    slots: list of dicts {label, name, size, where}
    images: list of (file_basename_without_ext, caption) — files under docs/screenshots/annotated/
    fallback_note: shown when images is empty/missing — explain why no screenshot.
    """
    s = prs.slides.add_slide(BLANK)
    add_bg(s); add_left_marker(s); add_brand_mark(s); add_page_no(s, n, total)
    add_eyebrow(s, label)
    add_h1(s, title)
    add_subtitle(s, lede, y=Inches(2.1))

    # Left column — slot list
    lx, ly = Inches(0.7), Inches(2.95)
    lw = Inches(4.6)
    add_text(s, lx, ly, lw, Inches(0.35),
             "📷  사진 슬롯 (admin 입력)", size=12, bold=True, color=TEXT, auto_grow=False)

    cy = ly + Inches(0.45)
    card_h = Inches(0.95)
    for slot in slots:
        add_rect(s, lx, cy, lw, card_h, WHITE, line=LINE)
        add_rect(s, lx, cy, Inches(0.05), card_h, GOLD)
        # circle badge (rendered as filled rect because pptx doesn't have round)
        badge_size = Inches(0.42)
        bx = lx + Inches(0.18)
        by = cy + (card_h - badge_size) / 2
        add_rect(s, bx, by, badge_size, badge_size, GOLD)
        add_text(s, bx, by + Inches(0.03), badge_size, Inches(0.36),
                 slot['label'], size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font=EN_FONT, auto_grow=False)
        # text block
        tx = bx + badge_size + Inches(0.18)
        tw = lw - (tx - lx) - Inches(0.15)
        add_text(s, tx, cy + Inches(0.12), tw, Inches(0.3),
                 slot['name'], size=12, bold=True, color=TEXT, auto_grow=False)
        add_text(s, tx, cy + Inches(0.4), tw, Inches(0.25),
                 slot['size'], size=10, color=GOLD, italic=True, auto_grow=False)
        add_text(s, tx, cy + Inches(0.62), tw, Inches(0.28),
                 slot['where'], size=10, color=TEXT_MUTED, auto_grow=False)
        cy += card_h + Inches(0.12)

    # Right column — image(s) or note
    rx, ry = Inches(5.55), Inches(2.95)
    rw = Inches(7.3)
    add_text(s, rx, ry, rw, Inches(0.35),
             "🌐  실제 사이트 화면 — 박스가 표시된 곳에 사진이 들어가요",
             size=12, bold=True, color=TEXT, auto_grow=False)

    avail_top = ry + Inches(0.45)
    avail_h   = Inches(3.6)   # leave room below for image bottom
    caption_h = Inches(0.25)
    if images:
        n = len(images)
        slot_total_h = avail_h / n
        slot_img_h   = slot_total_h - caption_h - Inches(0.08)
        for i, (basename, caption) in enumerate(images):
            img_path = os.path.join(SCREEN_DIR, f'{basename}.png')
            if not os.path.exists(img_path):
                continue
            from PIL import Image as PILImage
            with PILImage.open(img_path) as p:
                w_px, h_px = p.width, p.height
            ratio = w_px / h_px
            # fit width first
            target_w = rw
            target_h = target_w / ratio
            # cap height if it would exceed slot
            if target_h > slot_img_h:
                target_h = slot_img_h
                target_w = target_h * ratio
            ix = rx + (rw - target_w) / 2
            iy = avail_top + slot_total_h * i
            s.shapes.add_picture(img_path, ix, iy, width=target_w, height=target_h)
            if caption:
                cap_y = iy + target_h + Inches(0.04)
                add_text(s, rx, cap_y, rw, caption_h,
                         caption, size=9.5, color=TEXT_MUTED, italic=True,
                         align=PP_ALIGN.CENTER, auto_grow=False)
    elif fallback_note:
        add_rect(s, rx, avail_top, rw, Inches(2.4), BG_SOFT)
        add_text(s, rx + Inches(0.3), avail_top + Inches(0.3), rw - Inches(0.6), Inches(2.0),
                 fallback_note, size=12, color=TEXT_SOFT, line_spacing=1.55, auto_grow=False)
    return s

# ===============================================================
# Step slides
# ===============================================================
@reg
def step1(n, total):
    step_slide(n, total,
        step_label="STEP 1 · 카테고리 사진",
        step_title="카테고리 사진 채우기",
        est_time="약 15분",
        where="좌측 메뉴 → [카테고리]",
        click_path="카테고리 → 행 우측 ✏ → 사진 2장 업로드 → 저장",
        fields=[
            ("썸네일 사진",  "작은 카드용 · 약 800×800px",         True),
            ("히어로 사진",  "페이지 상단 큰 사진 · 약 1600×900px", True),
            ("기타 정보",    "이름은 자동 입력됨 — 손대지 마세요",   False),
        ],
        why="카테고리 8개 (얼굴/눈/코/바디/피부/모발/웰니스/치과) 는 자동으로 만들어져 있어요. 사진만 비어 있는 상태입니다. 사진을 넣으면 메인 화면 그리드에 즉시 반영됩니다.",
        sample=("얼굴 (face)",
                "썸네일: 얼굴 클로즈업\n히어로: 클리닉 분위기 사진"),
        warn=("카테고리는 새로 만들지 마세요",
              "8개가 사이트 분류 체계의 뼈대입니다. 추가/삭제하면 흔들립니다.",
              GOLD, GOLD_LIGHT, "✦"),
    )

@reg
def step1_usage(n, total):
    photo_usage_slide(n, total,
        label="STEP 1 · 카테고리 사진 — 사이트 어디에 보이나",
        title="카테고리 사진이 보이는 곳",
        lede="admin에서 올린 카테고리 사진은 홈의 'Browse by area' 그리드에 그대로 나옵니다.",
        slots=[
            {"label": "A", "name": "thumbnail_url",
             "size": "권장 800×800",
             "where": "BentoCategories 카드의 작은 사진 (fallback)"},
            {"label": "A", "name": "hero_image_url",
             "size": "권장 1600×900",
             "where": "카드의 메인 배경 사진 — 우선 사용"},
        ],
        images=[("01_home", "홈 페이지 → 'Browse by area' 섹션 (8 카드 모두 카테고리 hero/thumbnail)")],
    )

@reg
def step2(n, total):
    step_slide(n, total,
        step_label="STEP 2 · 고민 등록",
        step_title="고객의 고민(검색 키워드) 등록",
        est_time="약 30분",
        where="좌측 메뉴 → [고민]",
        click_path="고민 → [+ 새로 만들기] → 필드 입력 → 저장 → 반복",
        fields=[
            ("이름 한국어",  "예: 리프팅 / 여드름 흉터",     True),
            ("이름 영어",    "예: Lifting / Acne Scars",     True),
            ("이름 중국어",  "중국 시장 80% — 빠뜨리지 마세요", True),
            ("부위",        "얼굴/눈/코/바디/피부 중 하나",   True),
            ("URL 식별자",   "비워두면 자동 생성 · 건너뛰세요",  False),
        ],
        why="외국인 환자가 \"I want lifting\" 으로 검색하면, 우리 사이트에 'Lifting' 이라는 고민이 있어야 매칭됩니다. 이 단계가 매칭 결과의 입구를 만듭니다.",
        sample=("첫 10개 추천",
                "Lifting · Wrinkles · Sagging · Pores · Acne Scars · Dark Circles · Volume Loss · Nose Shape · Eye Shape"),
        warn=("목표 20~30개",
              "처음엔 인기 10개만 만들어도 사이트가 작동합니다. 나머지는 운영하며 추가.",
              GOLD, GOLD_LIGHT, "✦"),
    )

@reg
def step3(n, total):
    step_slide(n, total,
        step_label="STEP 3 · 시술 등록",
        step_title="시술 카탈로그 등록",
        est_time="약 2시간 (30~50개)",
        where="좌측 메뉴 → [시술]",
        click_path="시술 → [+ 새로 만들기] → 필드 입력 → 저장 → 반복",
        fields=[
            ("이름 한/영/중", "예: 울쎄라 / Ulthera Lifting",        True),
            ("카테고리",     "8개 중 선택 (콤보박스)",                True),
            ("기전",        "HIFU / RF / 필러 등 · 여러 개 가능",       True),
            ("강도",        "subtle / moderate / dramatic",          True),
            ("회복 기간",    "0~30일 · 시술 후 며칠 쉬어야 하나",         True),
            ("통증 강도",    "0~5 · 0=없음 / 5=매우 아픔",              True),
        ],
        why="사이트의 '메뉴판' 입니다. 시술이 없으면 병원이 어떤 시술을 등록할지 선택할 수 없어요. 인기 시술 30개를 먼저 만드세요.",
        sample=("우선 만들 30개",
                "HIFU · 써마지 · 필러 · 보톡스 · 리쥬란 · 실리프팅 · 코성형 · 쌍커풀 · 지방흡입 등"),
        warn=("한국어 + 영어 이름 모두 채우기",
              "외국인 환자는 영어, 운영자는 한국어로 봅니다.",
              GOLD, GOLD_LIGHT, "✦"),
    )

@reg
def step3_usage(n, total):
    photo_usage_slide(n, total,
        label="STEP 3 · 시술 사진 — 사이트 어디에 보이나",
        title="시술 사진이 보이는 곳",
        lede="시술 thumbnail 은 카드 형태로 거의 모든 페이지에서 노출됩니다. hero/gallery 는 시술 상세 페이지에서.",
        slots=[
            {"label": "A", "name": "thumbnail_url",
             "size": "권장 800×600",
             "where": "카테고리 그리드 / 매칭 결과 / AI 추천 카드 — 거의 모든 곳"},
            {"label": "B", "name": "thumbnail_url (또는 gallery_urls[0])",
             "size": "권장 1600×900",
             "where": "시술 상세 페이지 상단 hero 큰 사진"},
        ],
        images=[
            ("02_category", "카테고리 페이지 → 시술 그리드 카드 (A: thumbnail_url)"),
            ("03_treatment", "시술 상세 페이지 → 상단 hero (B: thumbnail_url 또는 gallery 첫 장)"),
        ],
    )

@reg
def step4(n, total):
    step_slide(n, total,
        step_label="STEP 4 · 매트릭스 (⭐가장 중요)",
        step_title="고민 ↔ 시술 연결하기",
        est_time="약 1~2시간",
        where="좌측 메뉴 → [고민↔시술 매트릭스]",
        click_path="고민 카드 안 [+ 시술 추가] → 시술 선택 → 관련도 지정",
        fields=[
            ("시술 선택",  "Step 3에서 만든 시술 중에서",          True),
            ("관련도",    "primary (1순위) / secondary / adjunct", True),
        ],
        why="매칭 알고리즘의 핵심 재료. 사이트 매칭 품질의 80%가 여기서 결정됩니다. 한 고민당 시술 2~3개씩 연결. 정성 들이는 만큼 결과가 좋아져요.",
        sample=("Lifting 고민 →",
                "HIFU (primary) · 실리프팅 (primary) · 써마지 (secondary) · 필러 (adjunct)"),
        warn=("운영의 진짜 가치는 여기",
              "어떤 고민에 어떤 시술이 좋은지 — 그 큐레이션이 우리의 차별점입니다.",
              GOLD, GOLD_LIGHT, "⭐"),
    )

# ===============================================================
# 10 STEP 5 개요
# ===============================================================
@reg
def step5_overview(n, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s); add_left_marker(s); add_brand_mark(s); add_page_no(s, n, total)
    add_eyebrow(s, "STEP 5 · 병원 등록 (가장 큰 작업)")
    add_meta_right(s, "예상 시간 · 병원당 약 30분")
    add_h1(s, "병원 한 곳을 등록하는 4단계")
    add_subtitle(s, "한 화면 안에서 4가지를 차례로 채웁니다.",
                 y=Inches(2.1))

    steps = [
        ("①", "기본 정보",      "병원 이름 · 주소 · 연락처 · 사진 · 응대 언어"),
        ("②", "의사 추가",      "병원 소속 의사 사진 + 직책 + 경력"),
        ("③", "시술 가격표",     "이 병원이 제공하는 시술 10~30개"),
        ("④", "전후사진",       "환자 동의 받은 before/after 사진"),
    ]
    cy = Inches(2.95)
    for num, title, body in steps:
        add_rect(s, Inches(0.7), cy, Inches(11.9), Inches(0.85), WHITE, line=LINE)
        add_rect(s, Inches(0.7), cy, Inches(0.08), Inches(0.85), GOLD)
        add_text(s, Inches(0.95), cy + Inches(0.2), Inches(0.7), Inches(0.5),
                 num, size=22, bold=True, color=GOLD, font=EN_FONT, auto_grow=False)
        add_text(s, Inches(1.65), cy + Inches(0.18), Inches(10.5), Inches(0.4),
                 title, size=15, bold=True, color=TEXT, auto_grow=False)
        add_text(s, Inches(1.65), cy + Inches(0.55), Inches(10.5), Inches(0.3),
                 body, size=11, color=TEXT_SOFT, auto_grow=False)
        cy += Inches(0.95)

    add_callout(s, Inches(0.7), Inches(6.4), Inches(11.9),
                "우선 7곳부터 → 전체 22곳",
                "Hershe · 노즈립 · 리엔장 강남 · 우아 · 365mc · Vellicell · Soi — 이 7곳 부터 시작하면 사이트가 살아납니다.",
                color=GOLD, bg=GOLD_LIGHT, icon="✦")

# ===============================================================
# STEP 5-1
# ===============================================================
@reg
def step5a(n, total):
    step_slide(n, total,
        step_label="STEP 5-1 · 병원 기본 정보",
        step_title="병원 기본 정보 입력",
        est_time="병원당 약 10분",
        where="좌측 메뉴 → [병원] → [+ 새로 만들기]",
        click_path="병원 → [+ 새로 만들기] → '기본 정보' 그룹부터",
        fields=[
            ("간판명 한국어/영어",   "예: 헤르쉬 청담점 / Hershe Cheongdam", True),
            ("도시 · 구 · 동",      "예: 서울 · 강남구 · 청담동",          True),
            ("응대 언어",          "한/영/중/일 등 · 여러 개 선택",         True),
            ("연락처",            "전화 / 카카오 / 위챗 / WhatsApp",     False),
            ("사진 (3종)",         "썸네일 · 히어로 · 갤러리",             True),
            ("계약 상태 = '활성' ⚠", "이거 안 바꾸면 사이트에 안 보임",      True),
        ],
        why="병원의 '대문' 정보. 사이트 병원 상세 페이지에 그대로 나타납니다. 외국인 환자가 가장 먼저 보는 화면이에요.",
        sample=("브랜드 정보 그룹",
                "같은 화면 아래쪽 '브랜드 정보' 그룹에 브랜드명·로고만 채우면 자동으로 브랜드가 등록됩니다."),
        warn=("계약 상태 = '활성(active)' 꼭 확인",
              "기본값은 '대기 중'. 활성으로 바꾸지 않으면 사이트에 절대 안 보입니다. 가장 흔한 실수.",),
    )

@reg
def step5a_usage(n, total):
    photo_usage_slide(n, total,
        label="STEP 5-1 · 병원 사진 — 사이트 어디에 보이나",
        title="병원 사진이 보이는 곳",
        lede="병원 hero 는 병원 상세 페이지 상단의 큰 사진. thumbnail 은 비교 카드에서.",
        slots=[
            {"label": "A", "name": "hero_image_url",
             "size": "권장 1600×900",
             "where": "병원 상세 페이지 상단의 큰 배경 사진"},
            {"label": "—", "name": "thumbnail_url",
             "size": "권장 800×600",
             "where": "시술 상세 페이지의 병원 비교 카드 · 디바이스 페이지"},
            {"label": "—", "name": "gallery_urls",
             "size": "권장 1200×800 × 5장",
             "where": "병원 상세 페이지 갤러리 섹션 (데이터 채워야 표시)"},
        ],
        images=[("04_hospital", "병원 상세 페이지 → 상단 hero (A: hero_image_url)")],
    )

# ===============================================================
# STEP 5-2
# ===============================================================
@reg
def step5b(n, total):
    step_slide(n, total,
        step_label="STEP 5-2 · 의사 추가",
        step_title="병원 소속 의사 추가",
        est_time="의사당 약 5분",
        where="병원 편집 화면 → '의사' 패널",
        click_path="병원 저장 후 → '의사' 패널 → [+ 의사 추가]",
        fields=[
            ("이름 한국어/영어",   "예: 김민준 / Dr. Min Jun Kim",  True),
            ("직책",            "예: 대표원장 / 부원장 / 전문의",     True),
            ("경력 한 줄",       "예: 강남대 출신, 15년 경력",       True),
            ("portrait 사진",   "정사각형 · 깔끔한 정면",           True),
            ("대표 의사 토글",    "병원의 얼굴이면 ON",              False),
        ],
        why="외국인 환자에게 의사 정보는 가장 큰 신뢰 시그널입니다. 사진과 경력만 있어도 결정률이 크게 올라가요.",
        sample=("처음엔",
                "병원당 대표 의사 1명만 먼저 등록. 부원장들은 나중에."),
        warn=("사진이 가장 중요",
              "이름·경력만 있고 사진이 없으면 신뢰도가 떨어져요. 병원에서 사진 받아오는 게 우선.",
              GOLD, GOLD_LIGHT, "✦"),
    )

@reg
def step5b_usage(n, total):
    photo_usage_slide(n, total,
        label="STEP 5-2 · 의사 사진 — 사이트 어디에 보이나",
        title="의사 사진이 보이는 곳",
        lede="병원 편집 화면에서 추가한 의사 사진은 병원 상세 페이지의 의사 카드 그리드에 나옵니다.",
        slots=[
            {"label": "—", "name": "portrait_url",
             "size": "권장 600×600 (정사각형)",
             "where": "병원 상세 페이지의 의사 카드 (프로필 사진)"},
            {"label": "—", "name": "hero_url",
             "size": "권장 1600×900",
             "where": "의사 상세 페이지의 hero (향후 페이지)"},
            {"label": "—", "name": "gallery_urls",
             "size": "권장 1200×800 × 여러 장",
             "where": "의사 상세 페이지의 갤러리 (향후)"},
        ],
        fallback_note=(
            "지금 사이트는 의사 데이터가 비어 있어서 화면 캡쳐가 없어요. "
            "병원 편집 페이지의 '의사' 패널에서 의사를 추가하시면 — 그 병원 상세 페이지에 "
            "자동으로 의사 카드가 나타납니다. "
            "포트레잇 사진 한 장만 있어도 외국인 환자에게 가장 큰 신뢰 시그널이 돼요."
        ),
    )

# ===============================================================
# STEP 5-3
# ===============================================================
@reg
def step5c(n, total):
    step_slide(n, total,
        step_label="STEP 5-3 · 시술 가격표",
        step_title="이 병원의 시술과 가격",
        est_time="병원당 10~30분",
        where="병원 편집 화면 → '시술 가격표' 패널",
        click_path="시술 패널 → [+ 시술 추가] → 시술 선택 → 가격 입력",
        fields=[
            ("시술 선택",       "Step 3에서 만든 시술 중에서",         True),
            ("시작 가격 (KRW)",  "예: 390000 (콤마 없이 숫자)",        True),
            ("원가 (선택)",     "할인 표시용 · 자동 % 계산",            False),
            ("장비 / 기기",     "예: Shurink, Thermage FLX",         False),
            ("시그너처 토글",    "대표 시술이면 ON",                  False),
            ("공개 여부",       "기본 ON · 가끔 숨길 때만 OFF",         True),
        ],
        why="병원 상세 페이지의 '메뉴판'. 환자가 어떤 시술을 얼마에 받을 수 있는지 결정해주는 핵심 데이터입니다.",
        sample=("한 병원당 보통",
                "10~30개 시술. 시그너처 1~3개는 꼭 토글 ON."),
        warn=("가격 받기 어려우면",
              "[공개 안 함] 토글로 사이트에 '상담문의'로 표시됩니다.",
              GOLD, GOLD_LIGHT, "✦"),
    )

# ===============================================================
# STEP 5-4 (special — B&A)
# ===============================================================
@reg
def step5d(n, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s); add_left_marker(s); add_brand_mark(s); add_page_no(s, n, total)
    add_eyebrow(s, "STEP 5-4 · 전후사진 — 가장 조심스럽게", color=WARN)
    add_meta_right(s, "사진당 약 5분")
    add_h1(s, "전후사진(B&A) 등록")

    add_callout(s, Inches(0.7), Inches(2.0), Inches(11.9),
                "이 단계는 법적으로 가장 민감합니다",
                "한국 PIPA 는 환자 사진을 민감정보로 분류. 환자 서면 동의 없는 사진은 절대 등록 X. 운영자 본인의 책임 영역입니다.",
                color=WARN, bg=WARN_BG)

    # Where strip
    add_rect(s, Inches(0.7), Inches(3.2), Inches(11.9), Inches(0.75), BG_SOFT)
    add_text(s, Inches(0.95), Inches(3.3), Inches(1.3), Inches(0.3),
             "📍  어디로", size=10, bold=True, color=GOLD, auto_grow=False)
    add_text(s, Inches(2.35), Inches(3.28), Inches(10.0), Inches(0.32),
             "병원 편집 화면 → 'B&A (전후사진)' 패널", size=12, color=TEXT,
             auto_grow=False)
    add_text(s, Inches(0.95), Inches(3.63), Inches(1.3), Inches(0.3),
             "👆  순서", size=10, bold=True, color=GOLD, auto_grow=False)
    add_text(s, Inches(2.35), Inches(3.63), Inches(10.0), Inches(0.32),
             "[+ B&A 추가] → 사진 2장 → 환자 메타 → 동의 체크 → 저장",
             size=11, color=TEXT_SOFT, italic=True, auto_grow=False)

    # Fields (left) + Tips (right)
    fields = [
        ("Before 사진 (시술 전)",            True),
        ("After 사진 (시술 후)",             True),
        ("이 사진의 시술",                   True),
        ("환자 나이대 · 환자 국가",           False),
        ("⚠ 환자 서면 동의 체크",            True),
        ("공개 범위 (전체 / 로그인 / 비공개)", True),
    ]
    add_text(s, Inches(0.7), Inches(4.2), Inches(7), Inches(0.35),
             "✏️  채워야 할 것", size=13, bold=True, color=TEXT, auto_grow=False)
    cy = Inches(4.6)
    for f, req in fields:
        add_text(s, Inches(0.7), cy, Inches(0.3), Inches(0.3),
                 "●" if req else "○",
                 size=13, color=(WARN if req else GOLD), auto_grow=False)
        add_text(s, Inches(1.0), cy, Inches(6.0), Inches(0.32),
                 f, size=11.5, color=TEXT, auto_grow=False)
        cy += Inches(0.36)

    add_text(s, Inches(7.95), Inches(4.2), Inches(4.7), Inches(0.35),
             "💡  팁", size=13, bold=True, color=TEXT, auto_grow=False)
    tips = [
        "처음엔 병원당 시그너처 시술 B&A 5장만",
        "동의서 없으면 일단 '비공개'로 등록",
        "환자 얼굴 식별 가능 사진은 더 조심",
        "의심스러우면 등록하지 마세요",
    ]
    cy = Inches(4.6)
    for t in tips:
        add_text(s, Inches(7.95), cy, Inches(0.25), Inches(0.32),
                 "·", size=14, color=GOLD, bold=True, auto_grow=False)
        add_text(s, Inches(8.2), cy, Inches(4.5), Inches(0.4),
                 t, size=11, color=TEXT_SOFT, auto_grow=False)
        cy += Inches(0.4)

@reg
def step5d_usage(n, total):
    photo_usage_slide(n, total,
        label="STEP 5-4 · 전후사진 — 사이트 어디에 보이나",
        title="B&A 사진이 보이는 곳",
        lede="병원 편집 화면에서 추가한 before/after 짝은 그 병원 상세 페이지의 B&A 섹션에 카드 형태로 나옵니다.",
        slots=[
            {"label": "—", "name": "before_url",
             "size": "권장 1000×1000 (정사각형 권장)",
             "where": "병원 상세 페이지 B&A 섹션 — 좌측 (시술 전)"},
            {"label": "—", "name": "after_url",
             "size": "권장 1000×1000 (정사각형 권장)",
             "where": "병원 상세 페이지 B&A 섹션 — 우측 (시술 후)"},
        ],
        fallback_note=(
            "지금 사이트는 B&A 데이터가 비어 있어서 화면 캡쳐가 없어요. "
            "병원 편집 페이지의 'B&A' 패널에서 추가한 전후사진은 → 자동으로 그 병원 상세 페이지의 B&A 섹션에 노출됩니다. "
            "⚠ 등록 전 반드시 환자 서면 동의를 확인하세요 (한국 개인정보보호법). "
            "동의가 의심스러우면 '비공개' 로 등록 후 동의 받고 공개 전환."
        ),
    )

# ===============================================================
# STEP 6 실시간 피드
# ===============================================================
@reg
def step6(n, total):
    step_slide(n, total,
        step_label="STEP 6 · 실시간 피드 시드",
        step_title="메인 화면 후기 부트스트랩",
        est_time="약 10분",
        where="좌측 메뉴 → [실시간 피드]",
        click_path="실시간 피드 → [+ 새로 만들기] → 5~10개",
        fields=[
            ("이니셜",        "이름 첫 글자만 · 예: M.",         True),
            ("국가 코드/명",   "예: SG · Singapore",             True),
            ("관련 시술",     "Step 3에서 만든 시술 중",          True),
            ("결과",         "matched / quoted / booked 등",    True),
            ("운영자 시드 ON", "운영자가 만든 표시",                True),
            ("공개 ON",       "사이트에 보이게",                  True),
        ],
        why="메인 화면 상단 '실시간으로 매칭되고 있어요' 띠와 'Recently matched' 섹션을 살리는 시드. 운영자가 5~10개 만들어두면 실 환자 옵트인 시 자동으로 늘어납니다.",
        sample=("좋은 분포",
                "다양한 국가 (싱가포르·중국·일본·인니·미국) + 다양한 시술 (얼굴·바디·치과)"),
        warn=("is_seed 토글 ON",
              "운영자가 만든 시드와 실 환자 옵트인 항목을 구분하는 토글입니다.",
              GOLD, GOLD_LIGHT, "✦"),
    )

# ===============================================================
# 일상 1 — 파트너 신청
# ===============================================================
@reg
def daily_partner(n, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s); add_left_marker(s); add_brand_mark(s); add_page_no(s, n, total)
    add_eyebrow(s, "일상 · 매일 체크")
    add_h1(s, "신규 파트너 병원 신청 검토")
    add_subtitle(s, "사이트 우측 하단 '파트너 등록' 으로 들어온 신청 처리.",
                 y=Inches(2.1))

    steps = [
        ("1", "WhatsApp 알림 확인",
         "신청 들어오면 운영자에게 자동 알림. 클릭 시 신청 요약 카드."),
        ("2", "[파트너 신청서] 메뉴 진입",
         "사이드바 두 번째 → 리스트 → 카드 클릭."),
        ("3", "내용 검토",
         "면허 · 안전 기록 · 외국인 응대 능력 · 중복 여부 확인."),
        ("4", "[승인] 또는 [거절]",
         "승인: 브랜드+병원 자동 등록 ('대기 중' 상태). 거절: 파일 보관함으로 이동."),
        ("5", "승인 병원 → '활성' 전환",
         "[병원] 메뉴에서 행 찾아 계약 상태 '활성'으로 + 사진/의사/시술 채우기."),
    ]
    cy = Inches(2.75)
    row_h = Inches(0.62)
    row_gap = Inches(0.04)
    for num, t, b in steps:
        add_rect(s, Inches(0.7), cy, Inches(11.9), row_h, WHITE, line=LINE)
        add_rect(s, Inches(0.9), cy + Inches(0.12), Inches(0.38), Inches(0.38), GOLD)
        add_text(s, Inches(0.9), cy + Inches(0.14), Inches(0.38), Inches(0.36),
                 num, size=15, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font=EN_FONT, auto_grow=False)
        add_text(s, Inches(1.45), cy + Inches(0.08), Inches(11.0), Inches(0.3),
                 t, size=12.5, bold=True, color=TEXT, auto_grow=False)
        add_text(s, Inches(1.45), cy + Inches(0.34), Inches(11.0), Inches(0.28),
                 b, size=10, color=TEXT_SOFT, auto_grow=False)
        cy += row_h + row_gap

    add_callout(s, Inches(0.7), Inches(6.35), Inches(11.9),
                "자동 등록 ≠ 자동 공개",
                "승인하면 DB엔 들어가지만 '대기 중' 이라 사이트엔 안 보임. 의도된 2단계 안전망입니다.",
                color=GOLD, bg=GOLD_LIGHT, icon="✦")

# ===============================================================
# 일상 2 — 환자 삭제 요청
# ===============================================================
@reg
def daily_rights(n, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s); add_left_marker(s); add_brand_mark(s); add_page_no(s, n, total)
    add_eyebrow(s, "일상 · 환자 권리 요청 대응")
    add_h1(s, "환자가 데이터 삭제를 요청할 때")
    add_subtitle(s, "PIPA(한국) + GDPR(EU) 30일 처리 의무. 매우 중요.",
                 y=Inches(2.1))

    steps = [
        ("1", "이메일 확인",
         "glowupinseoul@gmail.com 으로 \"erasure request\" 또는 \"개인정보 삭제 요청\" 수신."),
        ("2", "본인 확인",
         "해당 환자가 우리와 상담한 적 있는지 메일/WhatsApp 기록으로 확인."),
        ("3", "관련 데이터 찾기",
         "[실시간 피드] 메뉴에서 검색 → 공개 OFF 또는 행 삭제."),
        ("4", "스캔 기록 삭제",
         "별도 보관분은 개발자에게 위탁: \"세션 토큰 XXX 삭제\"."),
        ("5", "회신 + 기록 보존",
         "환자에게 \"30일 내 완료\" 회신. 회신 메일 캡쳐 보관 (분쟁 증거)."),
    ]
    cy = Inches(2.75)
    row_h = Inches(0.62)
    row_gap = Inches(0.04)
    for num, t, b in steps:
        add_rect(s, Inches(0.7), cy, Inches(11.9), row_h, WHITE, line=LINE)
        add_rect(s, Inches(0.9), cy + Inches(0.12), Inches(0.38), Inches(0.38), WARN)
        add_text(s, Inches(0.9), cy + Inches(0.14), Inches(0.38), Inches(0.36),
                 num, size=15, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font=EN_FONT, auto_grow=False)
        add_text(s, Inches(1.45), cy + Inches(0.08), Inches(11.0), Inches(0.3),
                 t, size=12.5, bold=True, color=TEXT, auto_grow=False)
        add_text(s, Inches(1.45), cy + Inches(0.34), Inches(11.0), Inches(0.28),
                 b, size=10, color=TEXT_SOFT, auto_grow=False)
        cy += row_h + row_gap

    add_callout(s, Inches(0.7), Inches(6.35), Inches(11.9),
                "30일 데드라인",
                "법적 의무 — 처리 누락 시 PIPA + GDPR 위반. 항상 회신 우선.",
                color=WARN, bg=WARN_BG)

# ===============================================================
# 막혔을 때
# ===============================================================
@reg
def help_slide(n, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s); add_left_marker(s); add_brand_mark(s); add_page_no(s, n, total)
    add_eyebrow(s, "막혔을 때")
    add_h1(s, "누구에게 어떻게 연락하나요?")

    cards = [
        ("Admin 화면이 안 열려요",
         "1. Ctrl+F5 새로고침\n2. 그래도 안 되면 개발자 연락",
         GOLD, GOLD_LIGHT),
        ("저장했는데 사이트에 안 보여요",
         "1. 계약 상태 '활성' 확인\n2. 사진 채웠는지 확인\n3. 새로고침",
         GOLD, GOLD_LIGHT),
        ("환자가 데이터 삭제 요청",
         "1. 30일 이내 처리 의무\n2. 직전 슬라이드 5단계대로",
         WARN, WARN_BG),
        ("환자가 부작용을 호소",
         "1. 즉시 해당 병원 연락\n2. 변호사 자문 검토\n3. 기록 남김",
         WARN, WARN_BG),
        ("법적인 질문",
         "1차: Privacy Policy 페이지\n2차: 변호사 상담",
         GOLD, GOLD_LIGHT),
        ("기능을 추가하고 싶어요",
         "운영자에게 한 줄로 요청\n개발자가 구현",
         GOLD, GOLD_LIGHT),
    ]
    cell_w = Inches(3.85); cell_h = Inches(1.7); gap = Inches(0.2)
    grid_x = Inches(0.7); grid_y = Inches(2.4)
    for i, (t, b, color, bg) in enumerate(cards):
        cx = grid_x + (cell_w + gap) * (i % 3)
        cy = grid_y + (cell_h + gap) * (i // 3)
        add_rect(s, cx, cy, cell_w, cell_h, bg)
        add_rect(s, cx, cy, Inches(0.06), cell_h, color)
        add_text(s, cx + Inches(0.25), cy + Inches(0.15), cell_w - Inches(0.4), Inches(0.4),
                 t, size=12, bold=True, color=TEXT, line_spacing=1.2, auto_grow=False)
        add_text(s, cx + Inches(0.25), cy + Inches(0.65), cell_w - Inches(0.4), Inches(1.0),
                 b, size=10.5, color=TEXT_SOFT, line_spacing=1.4, auto_grow=False)

    add_callout(s, Inches(0.7), Inches(6.45), Inches(11.9),
                "주요 연락처",
                "개발자(태희) WhatsApp · 운영자 슬랙/카톡 · 변호사(추후) · glowupinseoul@gmail.com",
                color=GOLD, bg=GOLD_LIGHT, icon="✦")

# ===============================================================
# 마무리
# ===============================================================
@reg
def closing(n, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s); add_left_marker(s); add_brand_mark(s); add_page_no(s, n, total)
    add_eyebrow(s, "마무리 · 일상 체크리스트")
    add_h1(s, "운영 리듬")
    add_subtitle(s, "처음 등록을 끝낸 후엔 이 리듬으로 유지하시면 됩니다.",
                 y=Inches(2.1))

    cols = [
        ("매일", [
            "□  [파트너 신청서] 새 항목 확인",
            "□  WhatsApp 새 메시지 응답",
            "□  새 옵트인 후기 확인",
        ]),
        ("매주", [
            "□  [실시간 피드] 오래된 시드 정리",
            "□  새 B&A 사진 동의 확인",
            "□  파트너 신청 검토 대기 정리",
        ]),
        ("매월", [
            "□  병원 가격·이벤트 최신화",
            "□  의사 변경 사항 반영",
            "□  Privacy Policy 변경 사항 검토",
            "□  새 시술 추가 (시장 반영)",
        ]),
    ]
    col_w = Inches(3.85); col_x0 = Inches(0.7); col_gap = Inches(0.2)
    for i, (title, items) in enumerate(cols):
        cx = col_x0 + (col_w + col_gap) * i
        add_rect(s, cx, Inches(2.95), col_w, Inches(3.7), WHITE, line=LINE)
        add_rect(s, cx, Inches(2.95), col_w, Inches(0.5), GOLD)
        add_text(s, cx, Inches(3.05), col_w, Inches(0.35),
                 title, size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, auto_grow=False)
        cy = Inches(3.65)
        for it in items:
            add_text(s, cx + Inches(0.25), cy, col_w - Inches(0.5), Inches(0.45),
                     it, size=11, color=TEXT, line_spacing=1.5, auto_grow=False)
            cy += Inches(0.5)

    add_text(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.4),
             "✦  Your Skin. Your Story. Seoul.  —  글로업 서울 운영팀",
             size=12, italic=True, color=GOLD,
             align=PP_ALIGN.CENTER, auto_grow=False)

# ===============================================================
def build():
    total = len(SLIDES)
    for i, fn in enumerate(SLIDES, 1):
        fn(i, total)
    base = os.path.normpath(os.path.join(os.path.dirname(__file__), "..", "docs"))
    out = os.path.join(base, "admin_guidebook.pptx")
    try:
        prs.save(out)
    except PermissionError:
        # File is locked (PowerPoint open?) — save with a versioned name instead.
        alt = os.path.join(base, "admin_guidebook_v2.pptx")
        prs.save(alt)
        out = alt
        print(f"[note] target file locked, wrote alternate: {alt}")
    print(f"[OK] saved: {out}")
    print(f"     slides: {total}")

if __name__ == "__main__":
    build()
